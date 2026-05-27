#!/usr/bin/env python3
"""Generate a modern presentation from a Halo API JSON BasicReport.

Produces a .pptx with native PowerPoint charts (no matplotlib):
  1. Campaign Overview — hero stats + publisher cards + insights
  2. Cross-Media Reach — column chart + metric cards
  3. Reach vs. Unique Reach — grouped column chart (if unique data)
  4. Frequency Distribution — column chart + per-publisher freq (if k+ data)
  5. Weekly Trends — line chart (if weekly data)
  6. Demographic Breakdown — table (if dimensionSpecSummary.groupings)
  7. Summary — stats + publisher table + takeaways

Usage:
    python3 generate_presentation.py <json_path> [output_path]
"""

import math
import os
import re
import sys
from dataclasses import dataclass, field
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SLIDE_WIDTH = Emu(9144000)
SLIDE_HEIGHT = Emu(5143500)


@dataclass
class BasicMetricSet:
    reach: int = 0
    percent_reach: float = 0.0
    k_plus_reach: list = field(default_factory=list)
    percent_k_plus_reach: list = field(default_factory=list)
    average_frequency: float = 0.0
    impressions: int = 0
    grps: float = 0.0


@dataclass
class ComponentMetrics:
    key: str = ""
    display_name: str = ""
    cumulative: Optional[BasicMetricSet] = None
    non_cumulative: Optional[BasicMetricSet] = None
    cumulative_unique_reach: int = 0
    non_cumulative_unique_reach: int = 0


@dataclass
class ComponentIntersection:
    component_keys: list = field(default_factory=list)
    cumulative_reach: int = 0


@dataclass
class ResultGroup:
    title: str = ""
    population_size: int = 0
    reporting_unit_cumulative: Optional[BasicMetricSet] = None
    reporting_unit_non_cumulative: Optional[BasicMetricSet] = None
    stacked_incremental_reach: list = field(default_factory=list)
    components: list = field(default_factory=list)
    component_intersections: list = field(default_factory=list)
    filter_name: str = ""
    dimension_groupings: list = field(default_factory=list)
    dimension_cell_label: str = ""
    metric_frequency: str = "total"
    metric_start_time: int = 0
    metric_end_time: int = 0


@dataclass
class BasicReport:
    name: str = ""
    title: str = ""
    campaign_group: str = ""
    campaign_group_display_name: str = ""
    start_date: str = ""
    end_date: str = ""
    duration_days: int = 0
    state: str = ""
    model_line: str = ""
    effective_model_line: str = ""
    iq_filters: list = field(default_factory=list)
    result_groups: list = field(default_factory=list)
    publisher_names: dict = field(default_factory=dict)


# ---------------------------------------------------------------------------
# JSON parser (Halo API responses)
# ---------------------------------------------------------------------------

# Known Origin/Aquila data provider display names.
# The Halo API often returns empty displayName fields for publishers;
# this map provides human-readable labels for known IDs.
KNOWN_DATA_PROVIDERS = {
    # Origin (UK)
    "dataProviders/F4iBUV7NkjQ": "Meta",
    "dataProviders/CqJcvwaa5tI": "TV",
    # Aquila (US)
    "dataProviders/YdyMFgVNkm0": "Meta",
    "dataProviders/bCx8khTNPmE": "TV",
}


def _json_date(d: dict) -> str:
    y = d.get("year")
    m = d.get("month")
    day = d.get("day")
    if y and m and day:
        return f"{int(y)}-{int(m):02d}-{int(day):02d}"
    return ""


def _json_epoch(val) -> int:
    if isinstance(val, dict):
        return int(val.get("seconds", 0) or 0)
    if isinstance(val, str) and val:
        from datetime import datetime

        try:
            dt = datetime.fromisoformat(val.replace("Z", "+00:00"))
            return int(dt.timestamp())
        except ValueError:
            return 0
    return 0


_AGE_LABELS = {
    "YEARS_16_TO_34": "16-34",
    "YEARS_18_TO_24": "18-24",
    "YEARS_18_TO_34": "18-34",
    "YEARS_25_TO_34": "25-34",
    "YEARS_35_TO_54": "35-54",
    "YEARS_55_PLUS": "55+",
    "YEARS_65_PLUS": "65+",
}


def _format_dimension_label(groupings: list) -> str:
    """Turn [{path:'common.sex', value:'MALE'}, {path:'common.age_group',
    value:'YEARS_55_PLUS'}] into 'Male 55+'. Order: sex first, then age."""
    sex = None
    age = None
    other = []
    for g in groupings:
        path = g.get("path", "")
        val = g.get("value", "")
        if path == "common.sex":
            sex = val.title() if val else None
        elif path == "common.age_group":
            age = _AGE_LABELS.get(val, val)
        else:
            other.append(f"{path.split('.')[-1]}={val}")
    parts = [p for p in [sex, age] if p] + other
    return " ".join(parts) if parts else ""


def _parse_json_metric_set(d: dict) -> BasicMetricSet:
    ms = BasicMetricSet()
    ms.reach = int(d.get("reach", 0) or 0)
    ms.percent_reach = float(d.get("percentReach", 0) or 0)
    ms.average_frequency = float(d.get("averageFrequency", 0) or 0)
    ms.impressions = int(d.get("impressions", 0) or 0)
    ms.grps = float(d.get("grps", 0) or 0)
    ms.k_plus_reach = [int(x) for x in (d.get("kPlusReach") or [])]
    ms.percent_k_plus_reach = [float(x) for x in (d.get("percentKPlusReach") or [])]
    return ms


def _parse_json_single_result(
    rg_title: str, res: dict, report: BasicReport
) -> ResultGroup:
    rg = ResultGroup()
    rg.title = rg_title

    ms_data = res.get("metricSet", {})
    rg.population_size = int(ms_data.get("populationSize", 0) or 0)

    meta = res.get("metadata", {})
    ru_summary = meta.get("reportingUnitSummary", {})
    for comp_summary in ru_summary.get("reportingUnitComponentSummary", []):
        key = comp_summary.get("component", "")
        name = comp_summary.get("displayName", "")
        if key and name:
            report.publisher_names[key] = name

    filt_meta = meta.get("filter", {})
    filt_iq = filt_meta.get("impressionQualificationFilter", "")
    if filt_iq:
        rg.filter_name = filt_iq

    freq_meta = meta.get("metricFrequency", {})
    if freq_meta.get("weekly"):
        rg.metric_frequency = "weekly"
    else:
        rg.metric_frequency = "total"

    ds_meta = meta.get("dimensionSpecSummary", {})
    groupings = ds_meta.get("groupings") or []
    for g in groupings:
        path = g.get("path", "")
        value = g.get("value", {})
        enum_val = value.get("enumValue", "") if isinstance(value, dict) else ""
        if path and enum_val:
            rg.dimension_groupings.append({"path": path, "value": enum_val})
    rg.dimension_cell_label = _format_dimension_label(rg.dimension_groupings)

    rg.metric_start_time = _json_epoch(meta.get("cumulativeMetricStartTime", {}))
    rg.metric_end_time = _json_epoch(meta.get("metricEndTime", {}))

    ru = ms_data.get("reportingUnit", {})
    cum = ru.get("cumulative")
    if cum:
        rg.reporting_unit_cumulative = _parse_json_metric_set(cum)
    nc = ru.get("nonCumulative")
    if nc:
        rg.reporting_unit_non_cumulative = _parse_json_metric_set(nc)

    sir = ru.get("stackedIncrementalReach", [])
    rg.stacked_incremental_reach = [int(x) for x in sir]

    for comp_data in ms_data.get("components", []):
        comp = ComponentMetrics()
        comp.key = comp_data.get("key", "")
        comp.display_name = (
            report.publisher_names.get(comp.key)
            or KNOWN_DATA_PROVIDERS.get(comp.key)
            or comp.key.split("/")[-1]
        )

        val = comp_data.get("value", {})
        c = val.get("cumulative")
        if c:
            comp.cumulative = _parse_json_metric_set(c)
        n = val.get("nonCumulative")
        if n:
            comp.non_cumulative = _parse_json_metric_set(n)

        cu = val.get("cumulativeUnique") or {}
        if cu:
            comp.cumulative_unique_reach = int(cu.get("reach", 0) or 0)
        ncu = val.get("nonCumulativeUnique") or {}
        if ncu:
            comp.non_cumulative_unique_reach = int(ncu.get("reach", 0) or 0)

        rg.components.append(comp)

    for ci_data in ms_data.get("componentIntersections", []):
        ci = ComponentIntersection()
        ci.component_keys = ci_data.get("components", [])
        ci_cum = ci_data.get("cumulative") or {}
        if ci_cum:
            ci.cumulative_reach = int(ci_cum.get("reach", 0) or 0)
        rg.component_intersections.append(ci)

    # Derive stacked incremental reach from components when not provided
    cum_total = rg.reporting_unit_cumulative
    if (
        not rg.stacked_incremental_reach
        and rg.components
        and cum_total
        and cum_total.reach > 0
    ):
        sorted_comps = sorted(
            rg.components,
            key=lambda c: c.cumulative.reach if c.cumulative else 0,
            reverse=True,
        )
        remaining = cum_total.reach
        for c in sorted_comps:
            cr = c.cumulative.reach if c.cumulative else 0
            incr = min(cr, remaining)
            rg.stacked_incremental_reach.append(incr)
            remaining -= incr
            if remaining <= 0:
                remaining = 0

    # Derive component intersections from reach overlap when not provided
    if (
        not rg.component_intersections
        and len(rg.components) == 2
        and cum_total
        and cum_total.reach > 0
    ):
        c0 = rg.components[0].cumulative
        c1 = rg.components[1].cumulative
        if c0 and c1:
            overlap = c0.reach + c1.reach - cum_total.reach
            if overlap > 0:
                ci = ComponentIntersection()
                ci.component_keys = [rg.components[0].key, rg.components[1].key]
                ci.cumulative_reach = overlap
                rg.component_intersections.append(ci)

    return rg


def _parse_json_result_group(rg_data: dict, report: BasicReport) -> list[ResultGroup]:
    title = rg_data.get("title", "")
    results = rg_data.get("results", [])
    if not results:
        return [ResultGroup(title=title)]

    # For total-frequency groups with a single result, return one ResultGroup.
    # For weekly groups (or groups with multiple results), expand each result
    # into its own ResultGroup so the rendering code sees the same structure
    # it gets from the API (one ResultGroup per week).
    first_meta = results[0].get("metadata", {})
    freq = first_meta.get("metricFrequency", {})
    is_weekly = bool(freq.get("weekly"))

    if not is_weekly and len(results) == 1:
        return [_parse_json_single_result(title, results[0], report)]

    groups = []
    for res in results:
        rg = _parse_json_single_result(title, res, report)
        groups.append(rg)
    return groups


def parse_json(filepath: str) -> BasicReport:
    import json

    with open(filepath) as f:
        data = json.load(f)

    report = BasicReport()
    report.name = data.get("name", "")
    report.title = data.get("title", "")
    report.campaign_group = data.get("campaignGroup", "")
    report.campaign_group_display_name = data.get("campaignGroupDisplayName", "")
    report.state = data.get("state", "")
    report.model_line = data.get("modelLine", "")
    report.effective_model_line = data.get("effectiveModelLine", "")

    ri = data.get("reportingInterval", {})
    rsd = ri.get("reportStartDate") or ri.get("effectiveReportStart", {})
    if rsd:
        report.start_date = _json_date(rsd)
    red = ri.get("reportEnd", {})
    if red:
        report.end_date = _json_date(red)

    if report.start_date and report.end_date:
        from datetime import date

        s = date.fromisoformat(report.start_date)
        e = date.fromisoformat(report.end_date)
        report.duration_days = (e - s).days

    for iqf in data.get("effectiveImpressionQualificationFilters", []):
        filt = iqf.get("impressionQualificationFilter", "")
        if filt:
            report.iq_filters.append(filt)
    report.iq_filters = list(dict.fromkeys(report.iq_filters))

    for rg_data in data.get("resultGroups", []):
        rgs = _parse_json_result_group(rg_data, report)
        report.result_groups.extend(rgs)

    return report


def parse_report(filepath: str) -> BasicReport:
    return parse_json(filepath)


def _enable_data_labels(
    series,
    num_format='#,##0,,"M"',
    position=None,
    font_size=Pt(8.5),
    bold=True,
    color=None,
):
    """Enable data labels with showVal=1 (python-pptx defaults to showVal=0)."""
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = font_size
    dl.font.bold = bold
    if color:
        dl.font.color.rgb = color
    dl.number_format = num_format
    if position is not None:
        try:
            dl.position = position
        except Exception:
            pass
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    show_val = dl._element.find(f"{{{ns_c}}}showVal")
    if show_val is not None:
        show_val.set("val", "1")
    else:
        sv = etree.SubElement(dl._element, f"{{{ns_c}}}showVal")
        sv.set("val", "1")
    return dl


# ── Formatting utilities ─────────────────────────────────────────────────────


def _fmt_num(n: int) -> str:
    if n >= 1_000_000:
        val = n / 1_000_000
        return f"{int(val)}M" if val == int(val) else f"{val:.1f}M"
    if n >= 1_000:
        val = n / 1_000
        return f"{int(val)}K" if val == int(val) else f"{val:.1f}K"
    return f"{n:,}"


def _fmt_full(n: int) -> str:
    return f"{n:,}"


def _fmt_pct(p: float) -> str:
    if p >= 1:
        return f"{p:.1f}%"
    return f"{p * 100:.1f}%"


def _short_name(display_name: str) -> str:
    replacements = [
        ("TV Broadcaster ", ""),
        ("National TV ", "National TV"),
        ("Digital Platform ", ""),
        ("Social Platform ", ""),
        (" Platform", ""),
        (" Broadcaster", ""),
    ]
    name = display_name
    for old, new in replacements:
        name = name.replace(old, new)
    return name.strip()


# ── Design tokens ────────────────────────────────────────────────────────────

# Typography
FONT = "Calibri Light"
FONT_BODY = "Calibri"

# Palette: neutral
INK = RGBColor(0x11, 0x18, 0x27)  # headings — almost black
TEXT = RGBColor(0x33, 0x3B, 0x4F)  # body text
TEXT_SEC = RGBColor(0x64, 0x6E, 0x83)  # secondary
TEXT_TER = RGBColor(0x94, 0x9C, 0xAE)  # labels, captions
DIVIDER = RGBColor(0xE2, 0xE5, 0xEB)  # lines
CARD_FILL = RGBColor(0xF8, 0xF9, 0xFB)  # card backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Palette: semantic
BLUE = RGBColor(0x2C, 0x5C, 0xE1)  # reach / primary accent
BLUE_LIGHT = RGBColor(0xEA, 0xF0, 0xFE)  # blue card bg
BLUE_DARK = RGBColor(0x1A, 0x3B, 0x8F)
GREEN = RGBColor(0x0D, 0x8A, 0x5E)  # incremental / growth
GREEN_LIGHT = RGBColor(0xE6, 0xF7, 0xF0)
VIOLET = RGBColor(0x6D, 0x3B, 0xD1)  # frequency
SLATE = RGBColor(0x47, 0x51, 0x69)  # impressions (neutral)

# Publisher colors — distinctive, accessible
PUB = [
    RGBColor(0x2C, 0x5C, 0xE1),  # blue
    RGBColor(0x8B, 0x5C, 0xF6),  # violet
    RGBColor(0x0D, 0x9B, 0x6A),  # teal
    RGBColor(0xE8, 0x6C, 0x00),  # amber
]

SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)

# ── Low-level helpers ────────────────────────────────────────────────────────


def _no_shadow(shape):
    """Strip all effects (shadow, glow, reflection) from a shape."""
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    # Explicit empty effectLst on the shape properties
    sp_pr = shape._element.find(f".//{ns_a}spPr")
    if sp_pr is not None:
        for el in sp_pr.findall(f"{ns_a}effectLst"):
            sp_pr.remove(el)
        for el in sp_pr.findall(f"{ns_a}effectDag"):
            sp_pr.remove(el)
        etree.SubElement(sp_pr, f"{ns_a}effectLst")
    # Neutralize theme effect reference (idx=0 means "no effect")
    style = shape._element.find(f"{ns_p}style")
    if style is not None:
        for ref in style.findall(f"{ns_a}effectRef"):
            ref.set("idx", "0")
            for child in list(ref):
                ref.remove(child)


def _rect(slide, left, top, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = line_w or Pt(0.5)
    else:
        s.line.fill.background()
    _no_shadow(s)
    return s


def _rounded(slide, left, top, w, h, fill=CARD_FILL, line=DIVIDER, radius=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    if s.adjustments and len(s.adjustments) > 0:
        s.adjustments[0] = radius
    _no_shadow(s)
    return s


def _dot(slide, left, top, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    _no_shadow(s)
    return s


def _txt(
    slide,
    left,
    top,
    w,
    h,
    text,
    size=Pt(10),
    bold=False,
    color=TEXT,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return tb


def _runs(slide, left, top, w, h, parts, align=PP_ALIGN.LEFT):
    """parts: [(text, size, bold, color, font?), ...]"""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for i, part in enumerate(parts):
        text, size, bold, color = part[:4]
        fn = part[4] if len(part) > 4 else FONT_BODY
        run = p.runs[0] if (i == 0 and p.runs) else p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = fn
    return tb


# ── Layout primitives ────────────────────────────────────────────────────────

MARGIN = Emu(530000)  # ~0.58" side margins
CONTENT_W = SLIDE_W - 2 * MARGIN


def _slide_header(slide, title, subtitle=None):
    """Thin accent bar + title + subtitle + divider."""
    _rect(slide, 0, 0, SLIDE_W, Emu(32000), fill=BLUE)

    _txt(
        slide,
        MARGIN,
        Emu(190000),
        CONTENT_W,
        Emu(380000),
        title,
        Pt(26),
        color=INK,
        font=FONT,
    )

    if subtitle:
        _txt(
            slide,
            MARGIN,
            Emu(560000),
            CONTENT_W,
            Emu(170000),
            subtitle,
            Pt(9),
            color=TEXT_SEC,
        )

    divider_y = Emu(780000) if subtitle else Emu(620000)
    _rect(slide, MARGIN, divider_y, CONTENT_W, Emu(8000), fill=DIVIDER)
    return divider_y + Emu(8000)


def _slide_footer(slide):
    _txt(
        slide,
        MARGIN,
        Emu(4950000),
        Emu(3500000),
        Emu(140000),
        "Source: Origin Cross-Media Measurement",
        Pt(6.5),
        color=TEXT_TER,
    )


# ── Stat card ────────────────────────────────────────────────────────────────


def _stat_card(
    slide, left, top, w, h, value, label, accent=BLUE, bg=CARD_FILL, border=DIVIDER
):
    _rounded(slide, left, top, w, h, fill=bg, line=border)
    _rect(slide, left, top + Emu(1), Emu(28000), h - Emu(2), fill=accent)

    _txt(
        slide,
        left + Emu(95000),
        top + Emu(50000),
        w - Emu(150000),
        Emu(130000),
        label,
        Pt(9),
        bold=True,
        color=TEXT_SEC,
    )
    _txt(
        slide,
        left + Emu(95000),
        top + Emu(210000),
        w - Emu(150000),
        Emu(230000),
        value,
        Pt(15),
        bold=True,
        color=accent,
        font=FONT,
    )


# ── Native PowerPoint chart ──────────────────────────────────────────────────


def _clean_chart_xml(chart):
    """Strip chrome from chart via XML: no border, transparent bg, subtle grid."""
    nsmap = {
        "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }
    root = chart._element

    # Remove chart-level border
    for sp in root.findall(".//c:spPr", nsmap):
        if sp.getparent().tag.endswith("}chartSpace") or sp.getparent() == root:
            ln = sp.find("a:ln", nsmap)
            if ln is None:
                ln = etree.SubElement(sp, f"{{{nsmap['a']}}}ln")
            for child in list(ln):
                ln.remove(child)
            etree.SubElement(ln, f"{{{nsmap['a']}}}noFill")

    # Style gridlines: very thin, light gray
    for gl in root.findall(".//c:majorGridlines", nsmap):
        sp = gl.find("c:spPr", nsmap)
        if sp is None:
            sp = etree.SubElement(gl, f"{{{nsmap['c']}}}spPr")
        ln = sp.find("a:ln", nsmap)
        if ln is None:
            ln = etree.SubElement(sp, f"{{{nsmap['a']}}}ln")
        ln.set("w", "6350")  # 0.5pt
        for child in list(ln):
            ln.remove(child)
        sf = etree.SubElement(ln, f"{{{nsmap['a']}}}solidFill")
        clr = etree.SubElement(sf, f"{{{nsmap['a']}}}srgbClr")
        clr.set("val", "ECEEF2")

    # Hide category axis line
    for catAx in root.findall(".//c:catAx", nsmap):
        sp = catAx.find("c:spPr", nsmap)
        if sp is None:
            sp = etree.SubElement(catAx, f"{{{nsmap['c']}}}spPr")
        ln = sp.find("a:ln", nsmap)
        if ln is None:
            ln = etree.SubElement(sp, f"{{{nsmap['a']}}}ln")
        for child in list(ln):
            ln.remove(child)
        sf = etree.SubElement(ln, f"{{{nsmap['a']}}}solidFill")
        clr = etree.SubElement(sf, f"{{{nsmap['a']}}}srgbClr")
        clr.set("val", "E2E5EB")


def _add_column_chart(
    slide, left, top, w, h, categories, values, colors, number_format='#,##0,,"M"'
):
    """Add a clean, modern column chart with per-bar colors."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Values", values)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, w, h, chart_data
    )
    chart = shape.chart
    chart.has_title = False
    chart.has_legend = False
    chart.style = 2

    plot = chart.plots[0]
    plot.gap_width = 100
    plot.overlap = 0

    series = plot.series[0]
    series.format.fill.background()

    # Color each bar individually
    for i, color in enumerate(colors):
        pt = series.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    # Data labels
    series.has_data_labels = True
    dl = series.data_labels
    dl.font.size = Pt(10)
    dl.font.bold = True
    dl.font.color.rgb = TEXT
    dl.font.name = FONT_BODY
    dl.number_format = number_format
    dl.show_value = True
    dl.show_category_name = False
    dl.show_series_name = False
    try:
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except:
        pass

    # Value axis: hidden labels, subtle gridlines
    vax = chart.value_axis
    vax.visible = False
    vax.has_major_gridlines = True
    vax.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xEE, 0xF2)
    vax.major_gridlines.format.line.width = Pt(0.5)

    # Category axis: clean, small labels
    cax = chart.category_axis
    cax.has_major_gridlines = False
    cax.tick_labels.font.size = Pt(8.5)
    cax.tick_labels.font.color.rgb = TEXT_SEC
    cax.tick_labels.font.name = FONT_BODY
    cax.format.line.color.rgb = RGBColor(0xE2, 0xE5, 0xEB)
    cax.format.line.width = Pt(0.5)

    _clean_chart_xml(chart)
    return shape


# ── Slide 1: Campaign Overview ───────────────────────────────────────────────


def build_slide_1(prs, report, rg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Subtitle line
    iq = []
    for f in report.iq_filters:
        if "mrc" in f.lower():
            iq.append("MRC")
        elif "ami" in f.lower():
            iq.append("AMI")
    sub = f"{report.start_date}  →  {report.end_date}  ·  {report.duration_days} days"
    if iq:
        sub += f"  ·  {' · '.join(iq)}"

    content_top = _slide_header(slide, report.campaign_group_display_name, sub)

    cum = rg.reporting_unit_cumulative
    pop = rg.population_size
    comps = rg.components

    # ── Stat cards: row of 4 ──
    card_w = Emu(1940000)
    card_h = Emu(520000)
    gap = Emu(90000)
    row_y = content_top + Emu(100000)

    stats = []
    if cum:
        if pop > 0:
            stats.append((_fmt_pct(cum.reach / pop * 100), "Reach", BLUE))
        if len(comps) >= 2 and rg.stacked_incremental_reach:
            incr = sum(rg.stacked_incremental_reach[1:])
            if cum.reach > 0:
                stats.append((f"+{_fmt_pct(incr / cum.reach)}", "Incremental", GREEN))
        if cum.average_frequency > 0:
            stats.append((f"{cum.average_frequency:.1f}x", "Frequency", VIOLET))
        if cum.impressions > 0:
            stats.append((_fmt_num(cum.impressions), "Impressions", SLATE))

    for i, (val, label, clr) in enumerate(stats[:4]):
        x = MARGIN + i * (card_w + gap)
        _stat_card(slide, x, row_y, card_w, card_h, val, label, accent=clr)

    # ── Publisher summary: compact horizontal cards ──
    pub_y = row_y + card_h + Emu(160000)
    _txt(
        slide,
        MARGIN,
        pub_y,
        CONTENT_W,
        Emu(180000),
        "Publisher contribution",
        Pt(9),
        bold=True,
        color=TEXT_SEC,
    )

    pub_card_y = pub_y + Emu(200000)
    total_reach = cum.reach if cum else 1
    n = len(comps)
    if n > 0:
        pub_w = (CONTENT_W - (n - 1) * Emu(70000)) // n
        for i, comp in enumerate(comps):
            x = MARGIN + i * (pub_w + Emu(70000))
            cr = comp.cumulative.reach if comp.cumulative else 0
            cf = comp.cumulative.average_frequency if comp.cumulative else 0
            share = cr / total_reach * 100 if total_reach else 0

            _rounded(
                slide, x, pub_card_y, pub_w, Emu(530000), fill=CARD_FILL, line=DIVIDER
            )
            # Color dot + name
            _dot(
                slide,
                x + Emu(80000),
                pub_card_y + Emu(70000),
                Emu(44000),
                PUB[i % len(PUB)],
            )
            _txt(
                slide,
                x + Emu(140000),
                pub_card_y + Emu(55000),
                pub_w - Emu(200000),
                Emu(130000),
                comp.display_name,
                Pt(8),
                bold=True,
                color=TEXT,
            )

            # Metrics row
            _runs(
                slide,
                x + Emu(80000),
                pub_card_y + Emu(230000),
                pub_w - Emu(160000),
                Emu(140000),
                [
                    (_fmt_num(cr), Pt(11), True, PUB[i % len(PUB)], FONT),
                    (f"  reach  ·  {share:.0f}%", Pt(8), False, TEXT_SEC),
                ],
            )

            _txt(
                slide,
                x + Emu(80000),
                pub_card_y + Emu(400000),
                pub_w - Emu(160000),
                Emu(110000),
                f"{cf:.1f}x freq  ·  {_fmt_num(comp.cumulative.impressions if comp.cumulative else 0)} imps",
                Pt(7.5),
                color=TEXT_TER,
            )

    # ── Key insights ──
    ins_y = pub_card_y + Emu(620000)
    _txt(
        slide,
        MARGIN,
        ins_y,
        CONTENT_W,
        Emu(170000),
        "Insights",
        Pt(9),
        bold=True,
        color=TEXT_SEC,
    )

    insights = []
    if cum and pop > 0:
        insights.append(
            f"Campaign reached {_fmt_full(cum.reach)} unique people "
            f"({_fmt_pct(cum.reach / pop * 100)} of universe) over {report.duration_days} days"
        )
    if len(comps) >= 2:
        anchor = comps[0]
        ar = anchor.cumulative.reach if anchor.cumulative else 0
        insights.append(
            f"{anchor.display_name} anchors with {_fmt_num(ar)} reach "
            f"({_fmt_pct(ar / cum.reach * 100 if cum and cum.reach else 0)} of total)"
        )
        for c in comps[1:]:
            if c.cumulative_unique_reach > 0:
                insights.append(
                    f"{c.display_name} adds {_fmt_num(c.cumulative_unique_reach)} "
                    f"unique people not reached elsewhere"
                )

    y = ins_y + Emu(200000)
    for t in insights[:3]:
        _runs(
            slide,
            MARGIN + Emu(30000),
            y,
            CONTENT_W - Emu(60000),
            Emu(180000),
            [("•  ", Pt(8), False, TEXT_TER), (t, Pt(8.5), False, TEXT)],
        )
        y += Emu(220000)

    _slide_footer(slide)


# ── Slide 2: Cross-Media Reach ───────────────────────────────────────────────


def build_slide_2(prs, report, rg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    cum = rg.reporting_unit_cumulative
    pop = rg.population_size
    total = cum.reach if cum else 0
    comps = rg.components
    n_pub = len(comps)

    sub = (
        f"Incremental reach across {n_pub} publishers"
        f"  ·  {_fmt_full(total)} unique people reached"
    )
    content_top = _slide_header(slide, "Cross-Media Reach", sub)

    # ── Left: Native column chart ──
    chart_left = MARGIN
    chart_top = content_top + Emu(70000)
    chart_w = Emu(5100000)
    chart_h = Emu(3500000)

    if rg.stacked_incremental_reach and comps:
        cats = []
        vals = rg.stacked_incremental_reach
        total_inc = sum(vals)
        for i, comp in enumerate(comps):
            name = _short_name(comp.display_name)
            pct = vals[i] / total_inc * 100 if total_inc else 0
            cats.append(f"{name}\n({pct:.0f}%)")

        colors = [PUB[i % len(PUB)] for i in range(len(comps))]
        _add_column_chart(
            slide, chart_left, chart_top, chart_w, chart_h, cats, vals, colors
        )

    # ── Right panel ──
    panel_left = chart_left + chart_w + Emu(160000)
    panel_w = SLIDE_W - panel_left - MARGIN

    # Card 1: Net Reach
    card_y = content_top + Emu(70000)
    card_h = Emu(700000)
    pad = Emu(130000)
    _rounded(
        slide,
        panel_left,
        card_y,
        panel_w,
        card_h,
        fill=BLUE_LIGHT,
        line=RGBColor(0xC7, 0xD7, 0xFE),
    )

    _txt(
        slide,
        panel_left + pad,
        card_y + Emu(55000),
        panel_w - pad * 2,
        Emu(120000),
        "Net Reach",
        Pt(8),
        bold=True,
        color=BLUE_DARK,
    )
    _txt(
        slide,
        panel_left + pad,
        card_y + Emu(200000),
        panel_w - pad * 2,
        Emu(280000),
        _fmt_pct(cum.percent_reach) if cum else "—",
        Pt(18),
        bold=True,
        color=BLUE,
        font=FONT,
    )
    _txt(
        slide,
        panel_left + pad,
        card_y + Emu(560000),
        panel_w - pad * 2,
        Emu(120000),
        f"{_fmt_full(total)} of {_fmt_full(pop)}",
        Pt(7.5),
        color=TEXT_SEC,
    )

    # Card 2: Incremental Reach
    c2_y = card_y
    if n_pub >= 2 and rg.stacked_incremental_reach:
        incr = sum(rg.stacked_incremental_reach[1:])
        incr_pct = incr / total * 100 if total else 0

        c2_y = card_y + card_h + Emu(90000)
        _rounded(
            slide,
            panel_left,
            c2_y,
            panel_w,
            card_h,
            fill=GREEN_LIGHT,
            line=RGBColor(0xA7, 0xF3, 0xD0),
        )

        _txt(
            slide,
            panel_left + pad,
            c2_y + Emu(55000),
            panel_w - pad * 2,
            Emu(120000),
            f"Incremental beyond {_short_name(comps[0].display_name)}",
            Pt(8),
            bold=True,
            color=RGBColor(0x06, 0x5F, 0x46),
        )
        _txt(
            slide,
            panel_left + pad,
            c2_y + Emu(200000),
            panel_w - pad * 2,
            Emu(280000),
            f"+{incr_pct:.1f}%",
            Pt(18),
            bold=True,
            color=GREEN,
            font=FONT,
        )
        _txt(
            slide,
            panel_left + pad,
            c2_y + Emu(560000),
            panel_w - pad * 2,
            Emu(120000),
            f"{_fmt_full(incr)} unique people",
            Pt(7.5),
            color=TEXT_SEC,
        )

    # Publisher legend
    legend_y = (
        c2_y + card_h + Emu(130000)
        if rg.stacked_incremental_reach
        else card_y + card_h + Emu(130000)
    )
    _txt(
        slide,
        panel_left,
        legend_y,
        panel_w,
        Emu(150000),
        "By publisher",
        Pt(8),
        bold=True,
        color=TEXT_SEC,
    )

    _rect(slide, panel_left, legend_y + Emu(170000), panel_w, Emu(6000), fill=DIVIDER)

    for i, comp in enumerate(comps):
        ry = legend_y + Emu(210000) + i * Emu(220000)
        cr = comp.cumulative.reach if comp.cumulative else 0
        cf = comp.cumulative.average_frequency if comp.cumulative else 0
        share = cr / total * 100 if total else 0

        _dot(
            slide,
            panel_left + Emu(10000),
            ry + Emu(30000),
            Emu(48000),
            PUB[i % len(PUB)],
        )
        _txt(
            slide,
            panel_left + Emu(80000),
            ry,
            Emu(1400000),
            Emu(120000),
            _short_name(comp.display_name),
            Pt(8),
            color=TEXT,
        )
        _txt(
            slide,
            panel_left + panel_w - Emu(1400000),
            ry,
            Emu(700000),
            Emu(120000),
            f"{_fmt_num(cr)} ({share:.0f}%)",
            Pt(8),
            bold=True,
            color=INK,
            align=PP_ALIGN.RIGHT,
        )
        _txt(
            slide,
            panel_left + panel_w - Emu(600000),
            ry,
            Emu(600000),
            Emu(120000),
            f"{cf:.1f}x",
            Pt(8),
            color=TEXT_TER,
            align=PP_ALIGN.RIGHT,
        )

    _slide_footer(slide)


# ── Slide 3: Reach vs. Unique Reach ──────────────────────────────────────────


def build_slide_reach_unique(prs, report, rg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    comps = rg.components
    cum = rg.reporting_unit_cumulative
    total = cum.reach if cum else 1

    content_top = _slide_header(
        slide,
        "Reach vs. Unique Reach",
        "Total reach includes audience overlap · Unique reach = people reached only by that publisher",
    )

    # Grouped bar chart: 2 series (Total Reach, Unique Reach)
    chart_data = CategoryChartData()
    cats = [_short_name(c.display_name) for c in comps]
    chart_data.categories = cats
    total_vals = [c.cumulative.reach if c.cumulative else 0 for c in comps]
    unique_vals = [c.cumulative_unique_reach for c in comps]
    chart_data.add_series("Total Reach", total_vals)
    chart_data.add_series("Unique Reach", unique_vals)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        MARGIN,
        content_top + Emu(80000),
        CONTENT_W,
        Emu(3400000),
        chart_data,
    )
    chart = shape.chart
    chart.has_title = False
    chart.style = 2

    plot = chart.plots[0]
    plot.gap_width = 80
    plot.overlap = -15

    # Series colors via XML to avoid python-pptx fill bugs
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    series_colors = [BLUE_LIGHT, BLUE]
    for idx, clr in enumerate(series_colors):
        ser = plot.series[idx]
        ser_el = ser._element
        sp_pr = ser_el.find(f"{{{ns_c}}}spPr")
        if sp_pr is None:
            sp_pr = etree.SubElement(ser_el, f"{{{ns_c}}}spPr")
        for old in sp_pr.findall(f"{{{ns_a}}}solidFill"):
            sp_pr.remove(old)
        sf = etree.SubElement(sp_pr, f"{{{ns_a}}}solidFill")
        srgb = etree.SubElement(sf, f"{{{ns_a}}}srgbClr")
        srgb.set("val", f"{clr[0]:02X}{clr[1]:02X}{clr[2]:02X}")

    # Data labels on both series
    for s in [plot.series[0], plot.series[1]]:
        dl = _enable_data_labels(s, position=XL_LABEL_POSITION.OUTSIDE_END, color=TEXT)
        dl.font.name = FONT_BODY

    # Legend
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.color.rgb = TEXT_SEC

    # Axes
    vax = chart.value_axis
    vax.visible = False
    vax.has_major_gridlines = True
    vax.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xEE, 0xF2)
    vax.major_gridlines.format.line.width = Pt(0.5)

    cax = chart.category_axis
    cax.tick_labels.font.size = Pt(8.5)
    cax.tick_labels.font.color.rgb = TEXT_SEC
    cax.format.line.color.rgb = DIVIDER
    cax.format.line.width = Pt(0.5)

    _clean_chart_xml(chart)

    # Key insight
    _rounded(
        slide,
        MARGIN,
        content_top + Emu(3650000),
        CONTENT_W,
        Emu(350000),
        fill=CARD_FILL,
        line=DIVIDER,
    )
    best = max(comps, key=lambda c: c.cumulative_unique_reach)
    _runs(
        slide,
        MARGIN + Emu(100000),
        content_top + Emu(3720000),
        CONTENT_W - Emu(200000),
        Emu(200000),
        [
            ("Key insight:  ", Pt(8), True, TEXT_SEC),
            (
                f"{best.display_name} delivers the most unique reach "
                f"({_fmt_num(best.cumulative_unique_reach)} people not reached elsewhere).",
                Pt(8.5),
                False,
                TEXT,
            ),
        ],
    )

    _slide_footer(slide)


# ── Slide 4: Frequency Distribution ──────────────────────────────────────────


def build_slide_frequency(prs, report, rg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cum = rg.reporting_unit_cumulative
    if not cum or not cum.k_plus_reach:
        return

    freq = cum.average_frequency
    content_top = _slide_header(
        slide,
        "Frequency Distribution",
        f"Average frequency: {freq:.1f}x  ·  K+ reach shows people reached at each frequency threshold",
    )

    # Column chart of k+ reach
    k_vals = cum.k_plus_reach
    cats = [f"{i + 1}+" for i in range(len(k_vals))]

    colors_freq = []
    for i in range(len(k_vals)):
        t = i / max(1, len(k_vals) - 1)
        r = int(44 + t * 100)
        g = int(92 + t * 20)
        b = int(225 - t * 100)
        colors_freq.append(RGBColor(r, g, b))

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series("K+ Reach", k_vals)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        MARGIN,
        content_top + Emu(80000),
        Emu(5600000),
        Emu(3400000),
        chart_data,
    )
    chart = shape.chart
    chart.has_title = False
    chart.has_legend = False
    chart.style = 2

    plot = chart.plots[0]
    plot.gap_width = 60
    series = plot.series[0]
    for i, clr in enumerate(colors_freq):
        series.points[i].format.fill.solid()
        series.points[i].format.fill.fore_color.rgb = clr

    dl = _enable_data_labels(
        series, font_size=Pt(8), position=XL_LABEL_POSITION.OUTSIDE_END, color=TEXT
    )
    dl.font.name = FONT_BODY

    vax = chart.value_axis
    vax.visible = False
    vax.has_major_gridlines = True
    vax.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xEE, 0xF2)
    cax = chart.category_axis
    cax.tick_labels.font.size = Pt(9)
    cax.tick_labels.font.color.rgb = TEXT_SEC
    cax.format.line.color.rgb = DIVIDER

    _clean_chart_xml(chart)

    # Right panel: per-publisher frequency
    panel_left = Emu(6100000)
    panel_w = SLIDE_W - panel_left - MARGIN
    comps = rg.components

    _txt(
        slide,
        panel_left,
        content_top + Emu(80000),
        panel_w,
        Emu(150000),
        "Per-publisher frequency",
        Pt(9),
        bold=True,
        color=TEXT_SEC,
    )

    _rect(
        slide, panel_left, content_top + Emu(260000), panel_w, Emu(6000), fill=DIVIDER
    )

    for i, comp in enumerate(comps):
        cf = comp.cumulative.average_frequency if comp.cumulative else 0
        ry = content_top + Emu(310000) + i * Emu(220000)

        _dot(
            slide,
            panel_left + Emu(10000),
            ry + Emu(30000),
            Emu(44000),
            PUB[i % len(PUB)],
        )
        _txt(
            slide,
            panel_left + Emu(70000),
            ry,
            Emu(1500000),
            Emu(120000),
            _short_name(comp.display_name),
            Pt(8),
            color=TEXT,
        )
        _txt(
            slide,
            panel_left + panel_w - Emu(600000),
            ry,
            Emu(600000),
            Emu(120000),
            f"{cf:.1f}x",
            Pt(10),
            bold=True,
            color=PUB[i % len(PUB)],
            align=PP_ALIGN.RIGHT,
        )

    # Effective frequency note
    pct_3plus = ""
    if len(k_vals) >= 3 and k_vals[0] > 0:
        pct_3plus = f"3+ reach is {k_vals[2] / k_vals[0] * 100:.0f}% of 1+ reach"

    if pct_3plus:
        note_y = content_top + Emu(310000) + len(comps) * Emu(220000) + Emu(120000)
        _rounded(
            slide,
            panel_left,
            note_y,
            panel_w,
            Emu(300000),
            fill=CARD_FILL,
            line=DIVIDER,
        )
        _txt(
            slide,
            panel_left + Emu(80000),
            note_y + Emu(40000),
            panel_w - Emu(160000),
            Emu(100000),
            "Effective frequency",
            Pt(7.5),
            bold=True,
            color=TEXT_SEC,
        )
        _txt(
            slide,
            panel_left + Emu(80000),
            note_y + Emu(160000),
            panel_w - Emu(160000),
            Emu(100000),
            pct_3plus,
            Pt(8),
            color=TEXT,
        )

    _slide_footer(slide)


# ── Slide: Demographic Breakdown ─────────────────────────────────────────────


def build_slide_demographics(prs, report, demo_rgs):
    """Render a sparse demographic table. Each row is one result-group cell;
    the segment label combines the RG title with the cell label (e.g.
    'Regular AMI · Male 55+'). Real reports have 1-N sparse cells, not a
    full demographic grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = _slide_header(
        slide, "Demographic Breakdown", "Per-segment reach and frequency"
    )

    col_widths = [
        Emu(2400000),  # Segment (wider — holds combined label)
        Emu(1100000),  # Population
        Emu(1100000),  # Reach
        Emu(1000000),  # % Reach
        Emu(1000000),  # Avg Freq
        Emu(1000000),  # Impressions
    ]
    headers = ["Segment", "Population", "Reach", "% Reach", "Avg Freq", "Impressions"]
    total_w = sum(col_widths)
    table_top = content_top + Emu(80000)

    _rounded(
        slide, MARGIN, table_top, total_w, Emu(260000),
        fill=RGBColor(0xF1, 0xF5, 0xF9), line=DIVIDER,
    )
    x = MARGIN
    for header, cw in zip(headers, col_widths):
        al = PP_ALIGN.LEFT if header == "Segment" else PP_ALIGN.RIGHT
        _txt(
            slide, x + Emu(40000), table_top + Emu(60000),
            cw - Emu(80000), Emu(130000),
            header, Pt(7.5), bold=True, color=TEXT_SEC, align=al,
        )
        x += cw

    for i, rg in enumerate(demo_rgs):
        row_y = table_top + Emu(290000) + i * Emu(250000)
        cum = rg.reporting_unit_cumulative
        if not cum:
            continue

        if i % 2 == 0:
            _rect(
                slide, MARGIN, row_y, total_w, Emu(240000),
                fill=RGBColor(0xFB, 0xFC, 0xFD),
            )
        _rect(slide, MARGIN, row_y + Emu(240000), total_w, Emu(4000), fill=DIVIDER)

        cell_label = rg.dimension_cell_label or "—"
        segment_text = (
            f"{rg.title} · {cell_label}" if rg.title else cell_label
        )

        vals = [
            segment_text,
            _fmt_full(rg.population_size),
            _fmt_full(cum.reach),
            _fmt_pct(cum.percent_reach),
            f"{cum.average_frequency:.1f}x" if cum.average_frequency else "—",
            _fmt_num(cum.impressions) if cum.impressions else "—",
        ]
        x = MARGIN
        for j, (val, cw) in enumerate(zip(vals, col_widths)):
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            clr = TEXT if j == 0 else INK
            _txt(
                slide, x + Emu(40000), row_y + Emu(55000),
                cw - Emu(80000), Emu(130000),
                val, Pt(8), bold=(j == 0), color=clr, align=al,
            )
            x += cw

    _slide_footer(slide)


# ── Slide 6: Summary ─────────────────────────────────────────────────────────


def build_slide_summary(prs, report, rg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cum = rg.reporting_unit_cumulative
    pop = rg.population_size
    comps = rg.components

    content_top = _slide_header(slide, "Summary")

    # Hero stat row (compact)
    stats_y = content_top + Emu(80000)
    card_w = Emu(2000000)
    gap = Emu(70000)

    metrics = []
    if cum:
        if pop > 0:
            metrics.append((_fmt_pct(cum.reach / pop * 100), "Net Reach", BLUE))
        if cum.average_frequency > 0:
            metrics.append((f"{cum.average_frequency:.1f}x", "Avg Frequency", VIOLET))
        if cum.impressions > 0:
            metrics.append((_fmt_num(cum.impressions), "Impressions", SLATE))
        if cum.grps > 0:
            metrics.append((f"{cum.grps:.1f}", "GRPs", GREEN))

    for i, (val, label, clr) in enumerate(metrics[:4]):
        x = MARGIN + i * (card_w + gap)
        _stat_card(slide, x, stats_y, card_w, Emu(500000), val, label, accent=clr)

    # Publisher share table
    table_y = stats_y + Emu(530000)
    _txt(
        slide,
        MARGIN,
        table_y,
        CONTENT_W,
        Emu(160000),
        "Campaign share by publisher",
        Pt(9),
        bold=True,
        color=TEXT_SEC,
    )

    col_widths = [
        Emu(2000000),
        Emu(1300000),
        Emu(1100000),
        Emu(1100000),
        Emu(1300000),
        Emu(1300000),
    ]
    headers = [
        "Publisher",
        "Reach",
        "% of Total",
        "Frequency",
        "Impressions",
        "Unique Reach",
    ]
    total_w = sum(col_widths)

    header_y = table_y + Emu(180000)
    _rounded(
        slide,
        MARGIN,
        header_y,
        total_w,
        Emu(240000),
        fill=RGBColor(0xF1, 0xF5, 0xF9),
        line=DIVIDER,
    )

    x = MARGIN
    for header, cw in zip(headers, col_widths):
        al = PP_ALIGN.LEFT if header == "Publisher" else PP_ALIGN.RIGHT
        _txt(
            slide,
            x + Emu(40000),
            header_y + Emu(55000),
            cw - Emu(80000),
            Emu(120000),
            header,
            Pt(7),
            bold=True,
            color=TEXT_SEC,
            align=al,
        )
        x += cw

    total_reach = cum.reach if cum else 1
    for i, comp in enumerate(comps):
        row_y = header_y + Emu(270000) + i * Emu(240000)
        cr = comp.cumulative.reach if comp.cumulative else 0
        cf = comp.cumulative.average_frequency if comp.cumulative else 0
        ci = comp.cumulative.impressions if comp.cumulative else 0
        cu = comp.cumulative_unique_reach
        share = cr / total_reach * 100 if total_reach else 0

        if i % 2 == 0:
            _rect(
                slide,
                MARGIN,
                row_y,
                total_w,
                Emu(230000),
                fill=RGBColor(0xFB, 0xFC, 0xFD),
            )
        _rect(slide, MARGIN, row_y + Emu(230000), total_w, Emu(3000), fill=DIVIDER)

        vals = [
            comp.display_name,
            _fmt_num(cr),
            f"{share:.1f}%",
            f"{cf:.1f}x",
            _fmt_num(ci),
            _fmt_num(cu) if cu else "—",
        ]
        x = MARGIN
        for j, (val, cw) in enumerate(zip(vals, col_widths)):
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            bld = j == 0
            clr = PUB[i % len(PUB)] if j == 0 else INK
            _txt(
                slide,
                x + Emu(40000),
                row_y + Emu(55000),
                cw - Emu(80000),
                Emu(120000),
                val,
                Pt(7.5),
                bold=bld,
                color=clr,
                align=al,
            )
            x += cw

    # Key takeaways
    take_y = header_y + Emu(270000) + len(comps) * Emu(240000) + Emu(140000)
    _txt(
        slide,
        MARGIN,
        take_y,
        CONTENT_W,
        Emu(150000),
        "Key takeaways",
        Pt(9),
        bold=True,
        color=TEXT_SEC,
    )

    points = []
    if cum and pop > 0:
        pct = cum.reach / pop * 100
        if pct > 50:
            points.append(
                f"Strong reach at {_fmt_pct(pct)} of universe ({_fmt_full(cum.reach)} people)"
            )
        else:
            points.append(
                f"Reach of {_fmt_pct(pct)} — room to grow with broader targeting or additional publishers"
            )

    if len(comps) >= 2 and rg.stacked_incremental_reach:
        incr = sum(rg.stacked_incremental_reach[1:])
        points.append(
            f"Cross-media strategy adds {_fmt_num(incr)} incremental people beyond the anchor publisher"
        )

    if cum and 3 <= cum.average_frequency <= 7:
        points.append(
            f"Healthy average frequency of {cum.average_frequency:.1f}x — within the 3–7x optimal range"
        )
    elif cum and cum.average_frequency > 7:
        points.append(
            f"Average frequency of {cum.average_frequency:.1f}x is elevated — review per-publisher frequency caps"
        )

    points.append(
        "Cost data not available in this report — overlay spend data for efficiency analysis"
    )

    y = take_y + Emu(180000)
    for t in points[:4]:
        _runs(
            slide,
            MARGIN + Emu(30000),
            y,
            CONTENT_W - Emu(60000),
            Emu(170000),
            [("•  ", Pt(7.5), False, TEXT_TER), (t, Pt(8), False, TEXT)],
        )
        y += Emu(200000)

    _slide_footer(slide)


# ── Main ─────────────────────────────────────────────────────────────────────


def _get_primary_rg(report):
    for r in report.result_groups:
        if not r.dimension_groupings and r.metric_frequency == "total" and r.components:
            if r.reporting_unit_cumulative and r.reporting_unit_cumulative.reach > 0:
                return r
    return report.result_groups[0] if report.result_groups else None


def _get_weekly_rgs(report):
    weekly = [
        rg
        for rg in report.result_groups
        if rg.metric_frequency == "weekly" and not rg.dimension_groupings
    ]
    if not weekly:
        return []
    # If weekly groups span multiple population sizes (e.g. IQF-filtered
    # variants), keep only those matching the largest so the chart plots
    # one consistent series instead of mixing scopes.
    max_pop = max(rg.population_size for rg in weekly)
    weekly = [rg for rg in weekly if rg.population_size == max_pop]
    weekly.sort(key=lambda rg: rg.metric_end_time)
    return weekly


def _get_demo_rgs(report):
    """Return total-frequency demographic cells only.

    Weekly demographic cells (one per week × cell) explode the row count and
    are a temporal breakdown, not a demographic summary — keep them out of
    this view so the table fits the slide. A typical real report has 3 RG
    titles (Regular/Incremental/Cumulative AMI) × 6 cells = 18 rows at most
    after this filter."""
    return [
        rg
        for rg in report.result_groups
        if rg.dimension_groupings
        and rg.reporting_unit_cumulative
        and rg.metric_frequency != "weekly"
    ]


def _build_slide_not_succeeded(prs, report):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=RGBColor(0xF8, 0xF9, 0xFB))
    _rect(slide, 0, 0, SLIDE_W, Emu(32000), fill=RGBColor(0xEF, 0x44, 0x44))

    msgs = {
        "RUNNING": "This report is still computing. Results are not yet available.",
        "FAILED": "This report failed to compute. Check the configuration and try again.",
        "INVALID": "This report has been invalidated. A new report may be needed.",
        "STATE_UNSPECIFIED": "This report has an unknown state.",
    }
    _txt(
        slide,
        Emu(500000),
        Emu(1800000),
        Emu(8000000),
        Emu(400000),
        report.title or "Untitled Report",
        Pt(22),
        bold=True,
        color=INK,
        font=FONT,
        align=PP_ALIGN.CENTER,
    )
    _txt(
        slide,
        Emu(500000),
        Emu(2400000),
        Emu(8000000),
        Emu(300000),
        f"Status: {report.state}",
        Pt(18),
        bold=True,
        color=RGBColor(0xEF, 0x44, 0x44),
        align=PP_ALIGN.CENTER,
    )
    _txt(
        slide,
        Emu(500000),
        Emu(2900000),
        Emu(8000000),
        Emu(300000),
        msgs.get(report.state, f"Report state: {report.state}"),
        Pt(11),
        color=TEXT_SEC,
        align=PP_ALIGN.CENTER,
    )


def build_slide_weekly(prs, report, weekly_rgs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    content_top = _slide_header(
        slide, "Weekly Reach Trends", "Cumulative reach build over campaign duration"
    )

    cats = []
    cum_vals = []
    nc_vals = []
    for i, rg in enumerate(weekly_rgs):
        if "(" in rg.title:
            label = rg.title.split("(")[-1].rstrip(")")
        elif rg.metric_end_time:
            from datetime import datetime, timezone

            dt = datetime.fromtimestamp(rg.metric_end_time, tz=timezone.utc)
            label = dt.strftime("%b %d")
        else:
            label = f"Wk {i + 1}"
        cats.append(label[:12])
        cum_vals.append(
            rg.reporting_unit_cumulative.reach if rg.reporting_unit_cumulative else 0
        )
        nc_vals.append(
            rg.reporting_unit_non_cumulative.reach
            if rg.reporting_unit_non_cumulative
            else 0
        )

    has_nc = any(v > 0 for v in nc_vals)

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series("Cumulative", cum_vals)
    if has_nc:
        chart_data.add_series("Weekly", nc_vals)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        MARGIN,
        content_top + Emu(80000),
        CONTENT_W,
        Emu(3400000),
        chart_data,
    )
    chart = shape.chart
    chart.has_title = False
    chart.style = 2

    s0 = chart.series[0]
    s0.format.line.color.rgb = BLUE
    s0.format.line.width = Pt(2.5)
    s0.smooth = True

    if has_nc:
        s1 = chart.series[1]
        s1.format.line.color.rgb = VIOLET
        s1.format.line.width = Pt(2)
        s1.format.line.dash_style = 4  # dash

    all_series = [s0] + ([chart.series[1]] if has_nc else [])
    for s in all_series:
        _enable_data_labels(
            s, font_size=Pt(8), position=XL_LABEL_POSITION.ABOVE, color=TEXT
        )

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.color.rgb = TEXT_SEC

    vax = chart.value_axis
    vax.visible = False
    vax.has_major_gridlines = True
    vax.major_gridlines.format.line.color.rgb = RGBColor(0xEC, 0xEE, 0xF2)
    cax = chart.category_axis
    cax.tick_labels.font.size = Pt(8)
    cax.tick_labels.font.color.rgb = TEXT_SEC
    cax.format.line.color.rgb = DIVIDER

    _clean_chart_xml(chart)

    # Insight
    if len(weekly_rgs) >= 2:
        fc = weekly_rgs[-1].reporting_unit_cumulative
        fi = weekly_rgs[0].reporting_unit_cumulative
        if fc and fi and fc.reach > 0:
            w1_share = fi.reach / fc.reach * 100
            _rounded(
                slide,
                MARGIN,
                content_top + Emu(3620000),
                CONTENT_W,
                Emu(300000),
                fill=CARD_FILL,
                line=DIVIDER,
            )
            _runs(
                slide,
                MARGIN + Emu(100000),
                content_top + Emu(3680000),
                CONTENT_W - Emu(200000),
                Emu(180000),
                [
                    ("Key insight:  ", Pt(8), True, TEXT_SEC),
                    (
                        f"Week 1 delivered {w1_share:.0f}% of final cumulative reach. "
                        f"Campaign reached {_fmt_num(fc.reach)} people over {len(weekly_rgs)} weeks.",
                        Pt(8.5),
                        False,
                        TEXT,
                    ),
                ],
            )

    _slide_footer(slide)


def generate_presentation(json_path: str, output_path: str = None) -> str:
    report = parse_report(json_path)

    if output_path is None:
        base = os.path.splitext(os.path.basename(json_path))[0]
        output_path = os.path.join(os.path.dirname(json_path), f"{base}.pptx")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    if report.state != "SUCCEEDED":
        _build_slide_not_succeeded(prs, report)
        prs.save(output_path)
        print(f"Generated: {output_path}")
        print(f"  Slides: {len(prs.slides)}")
        print(f"  State: {report.state}")
        return output_path

    rg = _get_primary_rg(report)
    if not rg:
        _build_slide_not_succeeded(prs, report)
        prs.save(output_path)
        return output_path

    weekly_rgs = _get_weekly_rgs(report)
    demo_rgs = _get_demo_rgs(report)

    if not rg.reporting_unit_cumulative and weekly_rgs:
        rg = weekly_rgs[-1]

    build_slide_1(prs, report, rg)

    if len(rg.components) >= 2:
        build_slide_2(prs, report, rg)

    has_unique = any(c.cumulative_unique_reach > 0 for c in rg.components)
    if len(rg.components) >= 2 and has_unique:
        build_slide_reach_unique(prs, report, rg)

    cum = rg.reporting_unit_cumulative
    if cum and cum.k_plus_reach:
        build_slide_frequency(prs, report, rg)

    if weekly_rgs:
        build_slide_weekly(prs, report, weekly_rgs)

    if demo_rgs:
        build_slide_demographics(prs, report, demo_rgs)

    build_slide_summary(prs, report, rg)

    prs.save(output_path)

    print(f"Generated: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Report: {report.title}")
    print(f"  State: {report.state}")
    if rg.components:
        print(f"  Publishers: {', '.join(c.display_name for c in rg.components)}")
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(f"Usage: {sys.argv[0]} <json_path> [output_path]")
        sys.exit(1)
    generate_presentation(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
