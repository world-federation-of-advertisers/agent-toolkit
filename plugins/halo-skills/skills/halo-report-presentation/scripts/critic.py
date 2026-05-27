#!/usr/bin/env python3.13
"""Auto-critic for generated XMM report presentations.

Evaluates each generated .pptx against quality criteria:
1. Data completeness — are all available metrics from the report present?
2. Slide count — appropriate number of slides for the data available?
3. Visual structure — headers, footers, consistent layout?
4. Text readability — no truncated numbers, proper formatting?
5. Chart presence — charts generated when data supports them?
6. Edge case handling — failed/running states, single publisher

Usage:
    python3.13 critic.py <json_dir> <output_dir>
"""

import os
import re
import sys
from dataclasses import dataclass, field

# Reuse the parser
sys.path.insert(0, os.path.dirname(__file__))
from generate_presentation import parse_report
from pptx import Presentation
from pptx.util import Emu


@dataclass
class CriticFinding:
    severity: str  # "CRITICAL", "WARNING", "INFO", "PASS"
    category: str
    message: str


@dataclass
class CriticReport:
    proto_file: str
    pptx_file: str
    findings: list = field(default_factory=list)
    score: float = 0.0

    def add(self, severity, category, message):
        self.findings.append(CriticFinding(severity, category, message))

    @property
    def critical_count(self):
        return sum(1 for f in self.findings if f.severity == "CRITICAL")

    @property
    def warning_count(self):
        return sum(1 for f in self.findings if f.severity == "WARNING")

    @property
    def pass_count(self):
        return sum(1 for f in self.findings if f.severity == "PASS")


def critique_presentation(proto_path: str, pptx_path: str) -> CriticReport:
    cr = CriticReport(proto_file=proto_path, pptx_file=pptx_path)
    report = parse_report(proto_path)

    if not os.path.exists(pptx_path):
        cr.add("CRITICAL", "file", f"Output file not found: {pptx_path}")
        return cr

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    slide_count = len(slides)

    # --- 1. State handling ---
    if report.state != "SUCCEEDED":
        if slide_count == 1:
            cr.add(
                "PASS",
                "state",
                f"Correctly produced single slide for {report.state} state",
            )
        else:
            cr.add(
                "CRITICAL",
                "state",
                f"Should produce 1 slide for {report.state} state, got {slide_count}",
            )

        # Check that the status message appears
        all_text = _extract_all_text(prs)
        if report.state in all_text:
            cr.add("PASS", "state", f"State '{report.state}' mentioned in slide text")
        else:
            cr.add(
                "WARNING", "state", f"State '{report.state}' not found in slide text"
            )
        cr.score = _calc_score(cr)
        return cr

    # --- 2. Slide count appropriateness ---
    primary_rg = None
    for rg in report.result_groups:
        if rg.reporting_unit_cumulative and rg.reporting_unit_cumulative.reach > 0:
            primary_rg = rg
            break
    if not primary_rg and report.result_groups:
        primary_rg = report.result_groups[0]

    if not primary_rg:
        cr.add("CRITICAL", "data", "No result groups found in parsed report")
        cr.score = _calc_score(cr)
        return cr

    n_publishers = len(primary_rg.components)
    has_weekly = any(rg.metric_frequency == "weekly" for rg in report.result_groups)
    has_kplus = (
        primary_rg.reporting_unit_cumulative
        and primary_rg.reporting_unit_cumulative.k_plus_reach
    )
    has_sir = bool(primary_rg.stacked_incremental_reach)

    expected_min = 3  # summary, total reach, final summary
    if n_publishers >= 2:
        expected_min += 2  # reach chart + unique reach
    if has_kplus:
        expected_min += 1  # frequency
    if has_weekly:
        expected_min += 1  # weekly

    if slide_count >= expected_min:
        cr.add(
            "PASS", "slide_count", f"{slide_count} slides (expected >= {expected_min})"
        )
    elif slide_count >= expected_min - 1:
        cr.add(
            "INFO",
            "slide_count",
            f"{slide_count} slides (expected >= {expected_min}, close)",
        )
    else:
        cr.add(
            "WARNING",
            "slide_count",
            f"{slide_count} slides (expected >= {expected_min})",
        )

    # --- 3. Data presence checks ---
    all_text = _extract_all_text(prs)
    cum = primary_rg.reporting_unit_cumulative

    # Campaign title
    if report.title and any(
        w.lower() in all_text.lower() for w in report.title.split()[:3]
    ):
        cr.add("PASS", "title", "Campaign title appears in presentation")
    else:
        cr.add(
            "WARNING", "title", f"Campaign title '{report.title}' not found in slides"
        )

    # Total reach number
    if cum:
        reach_str = f"{cum.reach:,}"
        reach_short = _fmt_num(cum.reach)
        if (
            reach_str in all_text
            or reach_short in all_text
            or str(cum.reach) in all_text
        ):
            cr.add("PASS", "reach", f"Total reach ({reach_str}) found in slides")
        else:
            cr.add("WARNING", "reach", f"Total reach ({reach_str}) not found in slides")

    # Publisher names
    for c in primary_rg.components:
        if c.display_name in all_text:
            cr.add("PASS", "publisher", f"Publisher '{c.display_name}' mentioned")
        else:
            cr.add(
                "WARNING",
                "publisher",
                f"Publisher '{c.display_name}' not found in slides",
            )

    # Frequency data
    if cum and cum.average_frequency > 0:
        freq_str = f"{cum.average_frequency:.1f}"
        if freq_str in all_text:
            cr.add("PASS", "frequency", f"Average frequency ({freq_str}) found")
        else:
            cr.add(
                "INFO", "frequency", f"Average frequency ({freq_str}) not found in text"
            )

    # K+ reach
    if has_kplus:
        has_freq_slide = any(
            "FREQUENCY" in _get_slide_text(s).upper()
            or "K+" in _get_slide_text(s).upper()
            for s in slides
        )
        if has_freq_slide:
            cr.add("PASS", "kplus", "Frequency distribution slide present")
        else:
            cr.add("WARNING", "kplus", "K+ reach data available but no frequency slide")

    # Weekly data
    if has_weekly:
        has_weekly_slide = any("WEEK" in _get_slide_text(s).upper() for s in slides)
        if has_weekly_slide:
            cr.add("PASS", "weekly", "Weekly trend slide present")
        else:
            cr.add("WARNING", "weekly", "Weekly data available but no weekly slide")

    # Incremental reach
    if has_sir and n_publishers >= 2:
        if "INCREMENTAL" in all_text.upper() or "incremental" in all_text.lower():
            cr.add("PASS", "incremental", "Incremental reach discussed")
        else:
            cr.add(
                "WARNING",
                "incremental",
                "Stacked incremental reach data available but not mentioned",
            )

    # --- 4. Visual structure checks ---
    for i, slide in enumerate(slides):
        text = _get_slide_text(slide)
        shapes = list(slide.shapes)

        # Check for header (first/second shape should be a header bar)
        if i == 0 or (slide_count > 1 and i > 0):
            has_header = any(
                s.shape_type == 1 and s.top < Emu(800000) and s.width > Emu(8000000)
                for s in shapes
            )
            if not has_header and len(text.strip()) > 10:
                cr.add(
                    "INFO",
                    "layout",
                    f"Slide {i + 1}: No full-width header bar detected",
                )

    # Check for source footer on data slides (not first/last)
    if slide_count > 2:
        for i in range(1, slide_count - 1):
            text = _get_slide_text(slides[i])
            if "Source:" in text or "Origin" in text:
                cr.add("PASS", "footer", f"Slide {i + 1}: Source attribution present")
                break
        else:
            cr.add("INFO", "footer", "No source attribution found on interior slides")

    # --- 5. Chart/image checks ---
    image_count = sum(
        1
        for slide in slides
        for shape in slide.shapes
        if shape.shape_type == 13  # PICTURE
    )
    if n_publishers >= 2 and image_count > 0:
        cr.add("PASS", "charts", f"{image_count} chart image(s) embedded")
    elif n_publishers >= 2 and image_count == 0:
        cr.add(
            "INFO", "charts", "No chart images found (matplotlib may not be available)"
        )

    # --- 6. Completeness of summary slide ---
    if slides:
        last_text = _get_slide_text(slides[-1]).upper()
        if (
            "SUMMARY" in last_text
            or "RECOMMEND" in last_text
            or "HIGHLIGHT" in last_text
        ):
            cr.add("PASS", "summary", "Final summary slide present")
        else:
            cr.add("WARNING", "summary", "No clear summary/recommendation slide at end")

    # Single publisher handling
    if n_publishers == 1:
        if slide_count >= 2:
            cr.add(
                "PASS",
                "single_pub",
                "Single publisher scenario produces reasonable deck",
            )
        else:
            cr.add(
                "WARNING",
                "single_pub",
                f"Single publisher scenario: only {slide_count} slide(s)",
            )

    # Campaign period check
    if report.start_date and report.end_date:
        if report.start_date in all_text or report.end_date in all_text:
            cr.add("PASS", "period", "Campaign period dates appear in slides")
        else:
            cr.add(
                "WARNING",
                "period",
                f"Campaign period ({report.start_date} to {report.end_date}) not shown",
            )

    # IQ filter context
    iq_mentioned = any(
        term in all_text.upper() for term in ["MRC", "AMI", "IMPRESSION QUALIFICATION"]
    )
    if report.iq_filters:
        if iq_mentioned:
            cr.add("PASS", "iq_filter", "IQ filter context mentioned in slides")
        else:
            cr.add("INFO", "iq_filter", "IQ filter type not mentioned in slides")

    # Per-publisher reach values should appear somewhere
    for c in primary_rg.components:
        if c.cumulative and c.cumulative.reach > 0:
            c_reach_str = f"{c.cumulative.reach:,}"
            if c_reach_str in all_text or str(c.cumulative.reach) in all_text:
                cr.add(
                    "PASS",
                    "pub_reach",
                    f"{c.display_name} reach ({c_reach_str}) in slides",
                )
            else:
                cr.add(
                    "INFO",
                    "pub_reach",
                    f"{c.display_name} reach ({c_reach_str}) not found in text",
                )

    cr.score = _calc_score(cr)
    return cr


def _extract_all_text(prs):
    texts = []
    for slide in prs.slides:
        texts.append(_get_slide_text(slide))
    return "\n".join(texts)


def _get_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text)
    return "\n".join(texts)


def _fmt_num(n):
    if n >= 1_000_000:
        return f"{n / 1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n / 1_000:.1f}K"
    return str(n)


def _calc_score(cr):
    if not cr.findings:
        return 0.0
    total = len(cr.findings)
    weights = {"PASS": 1.0, "INFO": 0.7, "WARNING": 0.3, "CRITICAL": 0.0}
    return sum(weights.get(f.severity, 0) for f in cr.findings) / total * 100


def run_batch_critique(proto_dir: str, output_dir: str):
    results = []
    protos = sorted(f for f in os.listdir(proto_dir) if f.endswith(".json"))

    for proto_file in protos:
        proto_path = os.path.join(proto_dir, proto_file)
        pptx_file = proto_file.replace(".json", ".pptx")
        pptx_path = os.path.join(output_dir, pptx_file)

        cr = critique_presentation(proto_path, pptx_path)
        results.append(cr)

    # Print report
    print("=" * 90)
    print(f"{'AUTO-CRITIC REPORT':^90}")
    print("=" * 90)
    print()

    total_score = 0
    for cr in results:
        proto_name = os.path.basename(cr.proto_file).replace(".json", "")
        status = (
            "PASS"
            if cr.critical_count == 0 and cr.warning_count == 0
            else "WARN"
            if cr.critical_count == 0
            else "FAIL"
        )
        icon = {"PASS": "✓", "WARN": "~", "FAIL": "✗"}[status]

        print(
            f"  {icon} {proto_name:45s} Score: {cr.score:5.1f}%  "
            f"(P:{cr.pass_count} W:{cr.warning_count} C:{cr.critical_count})"
        )

        total_score += cr.score

    avg_score = total_score / len(results) if results else 0
    print()
    print(f"{'─' * 90}")
    print(f"  Average Score: {avg_score:.1f}%")
    print(f"  Total Tests: {len(results)}")
    print(
        f"  Passed (no warnings): {sum(1 for r in results if r.warning_count == 0 and r.critical_count == 0)}"
    )
    print(
        f"  Warnings: {sum(1 for r in results if r.warning_count > 0 and r.critical_count == 0)}"
    )
    print(f"  Failed (critical): {sum(1 for r in results if r.critical_count > 0)}")
    print()

    # Detailed findings for non-PASS results
    print("=" * 90)
    print(f"{'DETAILED FINDINGS':^90}")
    print("=" * 90)
    print()

    for cr in results:
        non_pass = [f for f in cr.findings if f.severity != "PASS"]
        if non_pass:
            proto_name = os.path.basename(cr.proto_file).replace(".json", "")
            print(f"  {proto_name}:")
            for f in non_pass:
                icon = {"CRITICAL": "✗", "WARNING": "!", "INFO": "·"}[f.severity]
                print(f"    {icon} [{f.severity}] {f.category}: {f.message}")
            print()

    # Aggregate issue categories
    print("=" * 90)
    print(f"{'ISSUE SUMMARY BY CATEGORY':^90}")
    print("=" * 90)
    print()

    categories = {}
    for cr in results:
        for f in cr.findings:
            if f.severity in ("WARNING", "CRITICAL"):
                key = f.category
                if key not in categories:
                    categories[key] = []
                categories[key].append(f.message)

    if categories:
        for cat, msgs in sorted(categories.items(), key=lambda x: -len(x[1])):
            print(f"  {cat} ({len(msgs)} issues):")
            seen = set()
            for m in msgs:
                short = m[:100]
                if short not in seen:
                    seen.add(short)
                    print(f"    - {m[:120]}")
            print()
    else:
        print("  No warnings or critical issues found!")

    return results, avg_score


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <json_dir> <output_dir>")
        sys.exit(1)

    proto_dir = sys.argv[1]
    output_dir = sys.argv[2]
    results, avg_score = run_batch_critique(proto_dir, output_dir)
    sys.exit(0 if avg_score >= 70 else 1)
